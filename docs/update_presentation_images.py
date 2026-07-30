"""Add SWARM screenshots to industrial PPT and fix layout."""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "docs" / "presentation_assets"
PPT_DIR = Path(r"C:\Users\seena\Downloads")
PPT_NAME = "SWARM_Industrial_Presentation_TataSteel_2026"
PPT_SRC = PPT_DIR / f"{PPT_NAME}.pptx"
PPT_OUT = PPT_DIR / f"{PPT_NAME}_updated.pptx"
BACKUP = PPT_DIR / f"{PPT_NAME}.pptx.bak"
PRESENTER_LINE = "nanoFarm  •  July 2026"
DEMO_LINE = "Live demo: app.swarm.co.in/plant-hmi  |  tata.admin@tatasteel.com"

IMAGES = {
    "logo": ASSETS / "swarm-logo.png",
    "dashboard": ASSETS / "dashboard.png",
    "pid": ASSETS / "tata-steel-pid.png",
    "hmi": ASSETS / "plant-hmi-pfd.png",
    "connect": ASSETS / "connect-device.png",
}


def remove_shape(shape):
    el = shape.element
    el.getparent().remove(el)


def remove_shapes(slide, predicate):
    targets = [s for s in slide.shapes if predicate(s)]
    for shape in targets:
        remove_shape(shape)


def in_box(shape, left, top, width, height):
    return (
        shape.left >= left
        and shape.top >= top
        and shape.left <= left + width
        and shape.top <= top + height
    )


def add_picture_fit(slide, path, left, top, max_width, max_height):
    pic = slide.shapes.add_picture(str(path), left, top)
    ratio = pic.width / pic.height
    width = max_width
    height = int(width / ratio)
    if height > max_height:
        height = max_height
        width = int(height * ratio)
    pic.left = left + (max_width - width) // 2
    pic.top = top + (max_height - height) // 2
    pic.width = width
    pic.height = height
    return pic


def set_shape_text(shape, text, size=None):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if size:
        p.font.size = Pt(size)


def replace_text(slide, old, new):
    for shape in slide.shapes:
        if hasattr(shape, "text") and old in shape.text:
            shape.text = shape.text.replace(old, new)


def update_slide_1(slide):
    add_picture_fit(slide, IMAGES["logo"], Inches(0.55), Inches(0.35), Inches(1.35), Inches(1.05))
    replace_text(slide, "Presenter Name  •  Date", PRESENTER_LINE)


def update_slide_2(slide):
    remove_shapes(
        slide,
        lambda s: s.left >= Inches(6.4)
        and s.top >= Inches(1.7)
        and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE,
    )
    add_picture_fit(
        slide,
        IMAGES["dashboard"],
        Inches(6.45),
        Inches(1.85),
        Inches(6.55),
        Inches(4.65),
    )


def update_slide_3(slide):
    # Keep architecture diagram — no image overlay on this slide.
    pass


def update_slide_4(slide):
    remove_shapes(
        slide,
        lambda s: s.left >= Inches(4.7)
        and s.top >= Inches(2.9)
        and s.top <= Inches(5.2)
        and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE,
    )
    remove_shapes(
        slide,
        lambda s: hasattr(s, "text")
        and s.text
        and "Simple Process Flow" in s.text,
    )
    add_picture_fit(
        slide,
        IMAGES["pid"],
        Inches(4.75),
        Inches(2.95),
        Inches(8.35),
        Inches(4.15),
    )


def update_slide_5(slide):
    # Remove mock HMI blocks inside right panel
    remove_shapes(
        slide,
        lambda s: s.left >= Inches(5.0)
        and s.top >= Inches(1.8)
        and s.top <= Inches(6.6),
    )
    add_picture_fit(
        slide,
        IMAGES["hmi"],
        Inches(5.15),
        Inches(1.85),
        Inches(7.55),
        Inches(4.55),
    )
    # Pull status legend up (was off-slide)
    for shape in slide.shapes:
        if shape.top > Inches(5.95) and shape.left < Inches(2.5):
            shape.top = shape.top - Inches(1.35)


def update_slide_7(slide):
    remove_shapes(
        slide,
        lambda s: s.left >= Inches(5.0)
        and s.top >= Inches(1.8)
        and s.top <= Inches(5.0)
        and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE,
    )
    add_picture_fit(
        slide,
        IMAGES["dashboard"],
        Inches(5.1),
        Inches(1.85),
        Inches(7.65),
        Inches(3.05),
    )


def update_slide_8(slide):
    remove_shapes(
        slide,
        lambda s: s.left >= Inches(6.0)
        and s.top >= Inches(3.9)
        and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE,
    )
    remove_shapes(
        slide,
        lambda s: hasattr(s, "text")
        and s.text
        and "Same platform, single pilot" in s.text,
    )
    add_picture_fit(
        slide,
        IMAGES["connect"],
        Inches(6.05),
        Inches(3.95),
        Inches(7.55),
        Inches(2.85),
    )


def update_slide_10(slide):
    replace_text(
        slide,
        "Recommended next step: pilot deployment and live demonstration.",
        f"Recommended next step: pilot deployment and live demonstration.\n{DEMO_LINE}",
    )


def update_slide_9(slide):
    remove_shapes(
        slide,
        lambda s: s.left >= Inches(6.0)
        and s.top >= Inches(1.7)
        and s.top <= Inches(4.7)
        and s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE,
    )
    add_picture_fit(
        slide,
        IMAGES["dashboard"],
        Inches(6.15),
        Inches(1.82),
        Inches(7.35),
        Inches(2.95),
    )


def pick_source_ppt():
    if BACKUP.exists():
        try:
            prs = Presentation(str(BACKUP))
            if len(prs.slides[3].shapes) > 0:
                return BACKUP
        except Exception:
            pass
    if PPT_SRC.exists():
        return PPT_SRC
    raise FileNotFoundError(f"No presentation source found in {PPT_DIR}")


def main():
    for key, path in IMAGES.items():
        if not path.exists():
            raise FileNotFoundError(f"Missing asset {key}: {path}")

    source = pick_source_ppt()
    if not BACKUP.exists() and source != BACKUP:
        shutil.copy2(source, BACKUP)

    prs = Presentation(str(source))
    updaters = [
        update_slide_1,
        update_slide_2,
        None,
        update_slide_4,
        update_slide_5,
        None,
        update_slide_7,
        update_slide_8,
        update_slide_9,
        update_slide_10,
    ]
    for slide, updater in zip(prs.slides, updaters):
        if updater:
            updater(slide)

    if len(prs.slides[3].shapes) < 10:
        raise RuntimeError("Slide 4 is empty after update — aborting save")

    prs.save(str(PPT_OUT))
    docs_copy = ROOT / "docs" / PPT_OUT.name
    shutil.copy2(PPT_OUT, docs_copy)
    print(f"Source:  {source}")
    print(f"Updated: {PPT_OUT}")
    print(f"Docs:    {docs_copy}")
    print(f"Backup:  {BACKUP}")


if __name__ == "__main__":
    main()
